# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def main():
    out_path = r"C:\Users\Lenovo\Documents\kimi\workspace\天娱数科_深度分析_002354.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BLUE = RgbColor(0x1A, 0x23, 0x7E)
    LIGHT_BLUE = RgbColor(0x42, 0x5E, 0xB4)
    DARK_GRAY = RgbColor(0x2C, 0x3E, 0x50)
    WHITE = RgbColor(0xFF, 0xFF, 0xFF)
    GREEN = RgbColor(0x27, 0xAE, 0x60)
    RED = RgbColor(0xE7, 0x4C, 0x3C)

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BLUE
        shape.line.fill.background()
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = RgbColor(0xBB, 0xCC, 0xFF)
            p2.alignment = PP_ALIGN.CENTER
        return slide

    def add_content_slide(title, content_lines):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = DARK_BLUE
        bar.line.fill.background()
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.133), Inches(5.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(content_lines):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            if line.startswith("*") or line.startswith(">>") or line.startswith("<<") or line.startswith("#"):
                p.font.bold = True
        return slide

    def add_table_slide(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = DARK_BLUE
        bar.line.fill.background()
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        cols = len(headers)
        rows_count = len(rows) + 1
        table = slide.shapes.add_table(rows_count, cols, Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.6 * rows_count)).table
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(13)
                para.font.color.rgb = DARK_GRAY
                para.alignment = PP_ALIGN.CENTER
                if c_idx == 0:
                    para.alignment = PP_ALIGN.LEFT
        return slide

    # 1. Cover
    add_title_slide("天娱数科(002354.SZ)深度分析", "业务结构 | 客户分布 | 财务分析 | 2026半年报预测")

    # 2. Overview
    add_content_slide("一、公司概览", [
        "# 天娱数科(002354.SZ) 全称: 天娱数字科技(大连)集团股份有限公司",
        "# 所属行业: 传媒-数字营销 | 当前股价约6.57元 | 总市值约109亿元",
        "# 发展战略: 数字化、智能化、全球化",
        "# 2026年6月曾因概念炒作连续4个涨停，公司澄清无物理AI业务",
        "# 经营主题: 提质增效、强基拓新 - 推动AI能力平台化与全球化",
    ])

    # 3. Business structure
    add_content_slide("二、主营业务结构", [
        "# 1. AI营销SaaS (智能营销) - 核心增长引擎",
        "   - 自研天星大模型为底座，赋能智能营销",
        "   - 子公司椰子壳获评2026年Q2抖音电商金牌服务商",
        "   - 在马来西亚/印尼/泰国/日本取得TikTok服务商资质",
        "",
        "# 2. 移动应用分发PaaS",
        "   - 海外平台3uTools累计注册用户5775万人，月活232万人",
        "",
        "# 3. 空间智能MaaS",
        "   - 积累超150万条3D数据和65万条多模态数据",
        "   - Behavision空间智能大模型已完成备案",
        "",
        "# 4. 电竞游戏 (补充业务)",
        "   - Sunvy Poker流水同比增长31.8%，2026年2月登顶日本AppStore娱乐品类",
        "   - 飞升/苍穹变在东南亚保持良好运营",
    ])

    # 4. Revenue breakdown table
    add_table_slide("三、业务收入占比 (2025年半年报)", 
        ["业务板块", "收入(亿元)", "占比", "同比增速"],
        [
            ["数据流量行业(AI营销+应用分发)", "9.67", "97.93%", "+29.56%"],
            ["电竞游戏行业", "~0.20", "2.01%", "+57.04%"],
            ["合计", "9.88", "100%", "+29.64%"],
        ]
    )

    # 5. 2025 annual data
    add_table_slide("四、2025年全年业绩", 
        ["指标", "2025年", "2024年", "同比变化"],
        [
            ["营业收入", "20.77亿元", "15.79亿元", "+31.57%"],
            ["归母净利润", "-0.49亿元", "-1.18亿元", "减亏58.58%"],
            ["毛利额", "5.26亿元", "3.62亿元", "+45.40%"],
            ["毛利率", "25.31%", "22.90%", "+2.41pct"],
            ["海外收入", "0.57亿元", "0.09亿元", "+568%"],
            ["经营现金流净额", "1.16亿元", "0.32亿元", "+260.59%"],
        ]
    )

    # 6. Customer distribution
    add_content_slide("五、客户分布与海外市场", [
        ">> 国内客户 (AI营销):",
        "   3C数码: 小米、OPPO、TCL",
        "   家电: 海信、海尔、美的、石头科技",
        "   快消: 联合利华、几素",
        "   文娱: FIFA",
        "",
        ">> 海外布局 (东南亚为主):",
        "   马来西亚: 吉隆坡AI营销基地，TikTok Shop认证电商服务商",
        "   印尼/泰国/日本: TikTok服务商资质",
        "   Shopee直播核心合作伙伴",
        "   2025年海外收入5737万元，同比增长568%",
        "",
        ">> 用户数据:",
        "   移动应用分发平台3uTools: 累计注册用户3.49亿(2025年6月)",
        "   海外版3uTools: 累计注册用户5775万，月均活跃用户232万",
    ])

    # 7. 2025 Q3 report
    add_table_slide("六、2025年三季报核心数据", 
        ["指标", "2025年Q3累计", "同比变化", "2025年Q3单季"],
        [
            ["营业总收入", "15.08亿元", "+25.67%", "5.21亿元(+18.78%)"],
            ["归母净利润", "4270.31万元", "+597.6%", "1908.3万元(+1102.73%)"],
            ["毛利率", "24.33%", "+11.32pct", "-"],
            ["净利率", "3.18%", "+8213.61%", "-"],
            ["三费占营收比", "17.22%", "+1.91pct", "-"],
            ["每股经营性现金流", "0.05元", "+507.74%", "-"],
        ]
    )

    # 8. Q3 Pros
    add_content_slide("七、2025年三季报 - 优点", [
        "* 营收高增长: 前三季度营收15.08亿元，同比增长25.67%",
        "* 盈利能力大幅改善: 归母净利润4270万元，同比增长597.6%",
        "* 毛利率显著提升: 24.33%，同比提升11.32个百分点",
        "* 经营现金流改善: 每股经营性现金流0.05元，同比增长507.74%",
        "* 三费控制相对合理: 三费占营收比17.22%，保持可控",
        "* 数据流量业务持续放量: AI营销SaaS驱动核心增长",
    ])

    # 9. Q3 Cons
    add_content_slide("八、2025年三季报 - 缺点与风险", [
        "<< 历史盈利波动大: 近10年中位数ROIC为-3%，投资回报极差",
        "<< 每股净资产下降: 0.78元，同比减7.91%，股东权益有所侵蚀",
        "<< 商誉与减值压力: 2025年年报仍计提商誉减值4498万+长投减值2583万",
        "<< 主营业务造血能力待验证: 2024年亏损收窄主要源于减值计提减少",
        "<< 业务结构单一: 数据流量行业占比97.93%，对单一赛道依赖度高",
        "<< 现金流与利润背离: 2025年经营现金流净额1.16亿元，但净利润为-4862万元",
        "   差异系委托支付导致采购付现减少，并非真实经营改善",
    ])

    # 10. 2026 Q1 data
    add_table_slide("九、2026年一季报核心数据", 
        ["指标", "2026年Q1", "2025年Q1", "同比变化"],
        [
            ["营业收入", "5.47亿元", "4.85亿元", "+12.81%"],
            ["归母净利
